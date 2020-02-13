'''
Refer: https://github.com/Belval/pdf2image

1. Install poppler
Download http://blog.alivate.com.au/wp-content/uploads/2018/10/poppler-0.68.0_x86.7z
Extract to D:\RunNow\poppler-0.68.0\bin
Add to PATH: D:\RunNow\poppler-0.68.0\bin
'''


import os, sys

from PIL import Image
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

# Refer https://developers.google.com/slides/reference/rest/v1/Unit
# Declare size of A4 page by English Metric Unit (EMU)
A4_WIDTH = 9905695
A4_HEIGHT = 6858000

pdf_file = sys.argv[1]
print()
print("Converting file: " + pdf_file)
print()

# Prep presentation with a template size A4
prs = Presentation()

# Set slide dimensions A4 = 210x 297 mm
prs.slide_height = A4_HEIGHT
prs.slide_width = A4_WIDTH
    
blank_slide_layout = prs.slide_layouts[6]

# Create working folder
base_name = pdf_file.split(".pdf")[0]

# Convert PDF to list of images
print("Starting conversion...")
slideimgs = convert_from_path(pdf_file, 300, fmt= 'ppm', thread_count= 4)
print("...complete.")
print()

# Loop over slides
for i, slideimg in enumerate(slideimgs):
	if i % 10 == 0:
		print("Saving slide: " + str(i))

	imagefile = BytesIO()
	slideimg.save(imagefile, format='tiff')
	imagedata = imagefile.getvalue()
	imagefile.seek(0)
	width, height = slideimg.size

	# Add slide
	slide = prs.slides.add_slide(blank_slide_layout)
    
    # Scale image to fix with slide size A4.
    # Refer https://developers.google.com/slides/reference/rest/v1/Unit
	pic = slide.shapes.add_picture(imagefile, 0, 0, width= A4_WIDTH, height= A4_HEIGHT)

# Save Powerpoint
print()
print("Saving file: " + base_name + ".pptx")
prs.save(base_name + '.pptx')
print("Conversion complete. :)")
print()